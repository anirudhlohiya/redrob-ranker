import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Open presentation template
template_path = r"f:\redrob-ranker\Idea Submission Template _ Redrob.pptx"
output_path = r"f:\redrob-ranker\Idea Submission _ Redrob.pptx"

if not os.path.exists(template_path):
    print(f"Error: Template not found at {template_path}")
    exit(1)

prs = Presentation(template_path)

# Custom helper to style paragraphs & runs
def add_bullet(tf, text_runs, level=0):
    p = tf.add_paragraph()
    p.level = level
    # Add spacing between paragraphs
    p.space_after = Pt(6)
    
    for r_data in text_runs:
        r = p.add_run()
        r.text = r_data.get('text', '')
        if r_data.get('bold', False):
            r.font.bold = True
        r.font.name = 'Manrope'
        r.font.size = Pt(r_data.get('size', 13.5))
        r.font.color.rgb = RGBColor(32, 39, 41)
    return p

def format_slide_body(slide, body_shape, sections):
    tf = body_shape.text_frame
    tf.clear()
    tf.word_wrap = True
    
    is_first = True
    for sec in sections:
        # Title of section (e.g. bold header)
        if sec.get('header'):
            p = tf.paragraphs[0] if is_first else tf.add_paragraph()
            is_first = False
            p.space_after = Pt(8)
            p.space_before = Pt(8) if p != tf.paragraphs[0] else Pt(0)
            run = p.add_run()
            run.text = sec['header']
            run.font.bold = True
            run.font.name = 'Manrope'
            run.font.size = Pt(14.5)
            run.font.color.rgb = RGBColor(32, 39, 41)
            
        for bullet in sec.get('bullets', []):
            p = tf.add_paragraph()
            p.level = bullet.get('level', 0)
            p.space_after = Pt(4)
            # Add run items
            for run_data in bullet.get('runs', []):
                r = p.add_run()
                r.text = run_data.get('text', '')
                if run_data.get('bold', False):
                    r.font.bold = True
                r.font.name = 'Manrope'
                r.font.size = Pt(run_data.get('size', 12.5))
                r.font.color.rgb = RGBColor(32, 39, 41)

# Slide 1: Title Slide (index 0)
# Shape 1 (ID 55): 'Team Name :'
# Shape 2 (ID 56): 'Problem Statement :'
# Shape 3 (ID 57): 'Team Leader Name :'
slide1 = prs.slides[0]
for shape in slide1.shapes:
    if shape.has_text_frame:
        text = shape.text.strip()
        if "Team Name" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = "Team Name : "
            r.font.bold = True
            r.font.name = 'Manrope'
            r.font.color.rgb = RGBColor(32, 39, 41)
            r2 = p.add_run()
            r2.text = "Redrob-Ranker"
            r2.font.name = 'Manrope'
            r2.font.color.rgb = RGBColor(32, 39, 41)
        elif "Problem Statement" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = "Problem Statement : "
            r.font.bold = True
            r.font.name = 'Manrope'
            r.font.color.rgb = RGBColor(32, 39, 41)
            r2 = p.add_run()
            r2.text = "Build an offline, explainable, and CPU-efficient AI engine to rank candidate profiles against a 'Senior AI Engineer — Founding Team' role, filtering out data anomalies and keyword stuffers."
            r2.font.name = 'Manrope'
            r2.font.color.rgb = RGBColor(32, 39, 41)
        elif "Team Leader Name" in text:
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = "Team Leader Name : "
            r.font.bold = True
            r.font.name = 'Manrope'
            r.font.color.rgb = RGBColor(32, 39, 41)
            r2 = p.add_run()
            r2.text = "Anirudh Lohiya"
            r2.font.name = 'Manrope'
            r2.font.color.rgb = RGBColor(32, 39, 41)

# Slide 2: Solution Overview (index 1)
# Shape 2 (ID 64) is body
slide2 = prs.slides[1]
slide2_body = slide2.shapes[2]
slide2_content = [
    {
        'header': "Proposed Solution (Hybrid Multi-Stage Engine)",
        'bullets': [
            {
                'runs': [
                    {'text': "Stage 1 (Local Semantic Match): ", 'bold': True},
                    {'text': "Calculates semantic cosine similarity against the JD using localized TF-IDF (unigrams + bigrams) and Latent Semantic Analysis (LSA) projected to a 64D space, requiring zero remote APIs or large downloads."}
                ]
            },
            {
                'runs': [
                    {'text': "Stage 2 (Heuristic Rule Scorer): ", 'bold': True},
                    {'text': "Evaluates must-have competencies (vector search, ranking evaluation), nice-to-haves, preferred locations (Noida/Pune), and YOE bands (5–9 yrs) while applying soft penalties to title chasers or non-coding leads."}
                ]
            },
            {
                'runs': [
                    {'text': "Stage 3 (Behavioral & Integrity Gates): ", 'bold': True},
                    {'text': "Applies an availability multiplier (active recency, response rate, notice period) and intercepts inconsistent histories (overlapping employment, mismatched durations) via a Honeypot checker."}
                ]
            }
        ]
    },
    {
        'header': "Key Differentiators from Traditional Systems",
        'bullets': [
            {
                'runs': [
                    {'text': "Defeats Keyword Stuffing: ", 'bold': True},
                    {'text': "Analyzes complete career history and role tenures instead of simple phrase tallies, correctly downranking keyword-heavy but unqualified profiles (e.g. HR managers listing AI tools)."}
                ]
            },
            {
                'runs': [
                    {'text': "100% Deterministic & Safe: ", 'bold': True},
                    {'text': "Completely avoids hallucination risks, API latency, and data leakage by operating offline and locally without generative LLMs."}
                ]
            }
        ]
    }
]
format_slide_body(slide2, slide2_body, slide2_content)

# Slide 3: JD Understanding & Candidate Evaluation (index 2)
# Shape 2 (ID 71) is body
slide3 = prs.slides[2]
slide3_body = slide3.shapes[2]
slide3_content = [
    {
        'header': "Key Requirements Extracted from the JD",
        'bullets': [
            {
                'runs': [
                    {'text': "Production Embeddings & Retrieval: ", 'bold': True},
                    {'text': "Sentence-transformers, OpenAI embeddings, BGE, E5."}
                ]
            },
            {
                'runs': [
                    {'text': "Vector Search Infrastructure: ", 'bold': True},
                    {'text': "Pinecone, Weaviate, Qdrant, Milvus, OpenSearch, Elasticsearch, FAISS."}
                ]
            },
            {
                'runs': [
                    {'text': "Ranking Evaluation Frameworks: ", 'bold': True},
                    {'text': "NDCG, MRR, MAP, A/B testing, offline-to-online correlation."}
                ]
            },
            {
                'runs': [
                    {'text': "Ideal Profile Characteristics: ", 'bold': True},
                    {'text': "5–9 years of experience (ideal) and based in India (Pune/Noida preferred)."}
                ]
            }
        ]
    },
    {
        'header': "Evaluating Fit Beyond Keyword Matches",
        'bullets': [
            {
                'runs': [
                    {'text': "Disqualifier Analysis: ", 'bold': True},
                    {'text': "We model and filter out consulting-only, pure-research, and CV/Speech-only backgrounds without NLP exposure, as well as recent LangChain-hype profiles lacking pre-LLM ML depth."}
                ]
            },
            {
                'runs': [
                    {'text': "Stability & Role Checking: ", 'bold': True},
                    {'text': "Detects title-hopping (multiple short-tenure senior stints) and non-coding leads (roles with senior titles that lack hands-on evidence)."}
                ]
            }
        ]
    }
]
format_slide_body(slide3, slide3_body, slide3_content)

# Slide 4: Ranking Methodology (index 3)
# Shape 2 (ID 78) is body
slide4 = prs.slides[3]
slide4_body = slide4.shapes[2]
slide4_content = [
    {
        'header': "Stage-by-Stage Ranking Framework",
        'bullets': [
            {
                'runs': [
                    {'text': "1. Semantic Embedding Space (45% weight): ", 'bold': True},
                    {'text': "Fits TF-IDF Vectorizer (ngram_range=(1,2)) and TruncatedSVD (64 components) on candidate texts and JD. Computes cosine similarity to establish a normalized semantic overlap score."}
                ]
            },
            {
                'runs': [
                    {'text': "2. Heuristic Scoring (55% weight): ", 'bold': True},
                    {'text': "Awards points for production ML/retrieval keywords (max 40 pts), evaluation metrics (max 15 pts), Python (5 pts), nice-to-haves (10 pts), YOE compliance (15 pts), and location (10 pts). Applies penalties for title-chasers and non-coding leads (-8 pts each)."}
                ]
            },
            {
                'runs': [
                    {'text': "3. Hybrid Combination: ", 'bold': True},
                    {'text': "Aggregates scores into `Base Score = 0.55 * Rule_Score + 0.45 * Semantic_Score`. Hard-disqualified candidates are scaled by 0.15 to push them low while maintaining visibility."}
                ]
            },
            {
                'runs': [
                    {'text': "4. Behavioral Scaling: ", 'bold': True},
                    {'text': "Adjusts the base score with the `availability_multiplier` derived from inactivity duration, open-to-work flag, response rate, and notice period."}
                ]
            },
            {
                'runs': [
                    {'text': "5. Honeypot Penalty & Sorting: ", 'bold': True},
                    {'text': "Scales anomalous profiles by a 0.05x multiplier to push them to the absolute bottom of the pool. Sorts descending by score, using candidate ID ascending as a deterministic tie-breaker."}
                ]
            }
        ]
    }
]
format_slide_body(slide4, slide4_body, slide4_content)

# Slide 5: Explainability & Data Validation (index 4)
# Shape 2 (ID 85) is body
slide5 = prs.slides[4]
slide5_body = slide5.shapes[2]
slide5_content = [
    {
        'header': "Explainable Ranking Decisions",
        'bullets': [
            {
                'runs': [
                    {'text': "Deterministic Explanations: ", 'bold': True},
                    {'text': "Generates structured, clear reasons detailing years of experience, current titles, and exact keyword matches (e.g. `Senior AI Engineer with 5.9 yrs; 13 production ML/retrieval keywords; response rate 0.80`)."}
                ]
            },
            {
                'runs': [
                    {'text': "Eliminating Hallucinations: ", 'bold': True},
                    {'text': "Operating locally without a generative LLM means decisions are completely traceable, stable, and audit-friendly. The logic remains transparent during inspections."}
                ]
            }
        ]
    },
    {
        'header': "Data Integrity & Honeypot System",
        'bullets': [
            {
                'runs': [
                    {'text': "Chronology Checks: ", 'bold': True},
                    {'text': "Flags profiles containing overlapping employment intervals (>45 days across different companies) or inverted education dates."}
                ]
            },
            {
                'runs': [
                    {'text': "Implausible Durations: ", 'bold': True},
                    {'text': "Detects candidates claiming multiple concurrent current roles, total tenure exceeding stated YOE, or individual skill durations exceeding total career length."}
                ]
            },
            {
                'runs': [
                    {'text': "Soft-Gating Architecture: ", 'bold': True},
                    {'text': "Rather than deleting suspicious profiles (which creates false negatives), we scale their final scores by 0.05 to rank them at the bottom. This preserves data transparency for human review."}
                ]
            }
        ]
    }
]
format_slide_body(slide5, slide5_body, slide5_content)

# Slide 6: End-to-End Workflow (index 5)
# Shape 2 (ID 92) is body
slide6 = prs.slides[5]
slide6_body = slide6.shapes[2]
slide6_content = [
    {
        'header': "Core Processing Workflow",
        'bullets': [
            {
                'runs': [
                    {'text': "1. Candidate Ingestion: ", 'bold': True},
                    {'text': "Stream records from candidate JSONL file. Concatenate profile text fields (summary, headline, current title, job history titles/descriptions, skills) into a single corpus document."}
                ]
            },
            {
                'runs': [
                    {'text': "2. Semantic Vector Space: ", 'bold': True},
                    {'text': "Initialize and fit TF-IDF vectorizer (using unigrams and bigrams, up to 8000 features). Apply Truncated SVD to project the sparse matrix into a dense 64D space, and calculate cosine similarity against the parsed JD vector."}
                ]
            },
            {
                'runs': [
                    {'text': "3. Heuristic Scoring: ", 'bold': True},
                    {'text': "Extract candidate attributes and check against the JD rubric to calculate a rule-based fit score. Identify any hard disqualifiers (e.g. consulting-only) and apply a 0.15x scaling multiplier if present."}
                ]
            },
            {
                'runs': [
                    {'text': "4. Behavioral & Honeypot Processing: ", 'bold': True},
                    {'text': "Calculate availability multiplier using engagement signals. Scan dates for timeline overlaps or implausible durations, applying a 0.05x honeypot penalty to anomalous records."}
                ]
            },
            {
                'runs': [
                    {'text': "5. Ranking Blending & Output: ", 'bold': True},
                    {'text': "Compute final hybrid score: `(0.55 * Rule + 0.45 * Semantic) * Availability * Honeypot`. Sort descending, resolve ties by candidate ID, filter the top 100, and output to a validated CSV file."}
                ]
            }
        ]
    }
]
format_slide_body(slide6, slide6_body, slide6_content)

# Slide 7: System Architecture (index 6)
# Slide 7 has only a title, we need to add a body text box.
slide7 = prs.slides[6]
left = Inches(0.53)
top = Inches(1.50)
width = Inches(8.9)
height = Inches(3.5)
slide7_body = slide7.shapes.add_textbox(left, top, width, height)
slide7_content = [
    {
        'header': "Modular Architecture (Offline-First)",
        'bullets': [
            {
                'runs': [
                    {'text': "jd_config.py (Rubric Central): ", 'bold': True},
                    {'text': "Defines the core matching keywords, locations (Pune/Noida), target experience limits (5–9 yrs), and disqualifier indicators used across the pipeline."}
                ]
            },
            {
                'runs': [
                    {'text': "features.py (Feature Extractor): ", 'bold': True},
                    {'text': "Parses dates, calculates tenure lengths, extracts locations, evaluates behavioral factors, and runs the 5 logical timeline honeypot rules."}
                ]
            },
            {
                'runs': [
                    {'text': "scorer.py (Evaluation Blender): ", 'bold': True},
                    {'text': "Computes TF-IDF & Truncated SVD, measures cosine similarities, evaluates the rubric checklist, and blends results into final candidate scores."}
                ]
            },
            {
                'runs': [
                    {'text': "rank.py (Driver & CLI): ", 'bold': True},
                    {'text': "Streams records, coordinates feature extraction, scores candidate documents, performs sorting, and writes formatted CSV outputs."}
                ]
            }
        ]
    },
    {
        'header': "Architecture Highlights",
        'bullets': [
            {
                'runs': [
                    {'text': "High Scalability: ", 'bold': True},
                    {'text': "Processes candidate profiles dynamically to minimize system memory footprint, avoiding large array loads."}
                ]
            },
            {
                'runs': [
                    {'text': "No Network / GPU Required: ", 'bold': True},
                    {'text': "Avoids deep-learning model downloads, providing fast execution on minimal CPU hardware with standard Python packages."}
                ]
            }
        ]
    }
]
format_slide_body(slide7, slide7_body, slide7_content)

# Slide 8: Results & Performance (index 7)
# Shape 2 (ID 105) is body
slide8 = prs.slides[7]
slide8_body = slide8.shapes[2]
slide8_content = [
    {
        'header': "Runtime Efficiency on 100K Profile Pool",
        'bullets': [
            {
                'runs': [
                    {'text': "Execution Time: ", 'bold': True},
                    {'text': "Completed in ~100 seconds, significantly below the 5-minute challenge limit."}
                ]
            },
            {
                'runs': [
                    {'text': "Memory Usage: ", 'bold': True},
                    {'text': "Peak RAM stayed at ~2.8 GB, well within the 16 GB system memory allocation."}
                ]
            }
        ]
    },
    {
        'header': "Ranking Accuracy & Heuristic Validation",
        'bullets': [
            {
                'runs': [
                    {'text': "Sample Candidate Evaluation: ", 'bold': True},
                    {'text': "Tested against the 50-candidate sample. The ranker successfully placed the single candidate possessing production recommendation-system experience as #1 (score: 1.0)."}
                ]
            },
            {
                'runs': [
                    {'text': "Filtering Keyword Stuffers: ", 'bold': True},
                    {'text': "Unqualified candidates who stuffed AI terms into their profile summaries (e.g. non-technical managers) were correctly penalized and ranked near the bottom, avoiding keyword-matching pitfalls."}
                ]
            },
            {
                'runs': [
                    {'text': "Honeypot Gate In Action: ", 'bold': True},
                    {'text': "Profiles containing contradictory dates or education histories were successfully flagged and demoted, preserving output integrity."}
                ]
            }
        ]
    }
]
format_slide_body(slide8, slide8_body, slide8_content)

# Slide 9: Technologies Used (index 8)
# Shape 2 (ID 112) is body
slide9 = prs.slides[8]
slide9_body = slide9.shapes[2]
slide9_content = [
    {
        'header': "Selected Tech Stack & Rationale",
        'bullets': [
            {
                'runs': [
                    {'text': "Python 3: ", 'bold': True},
                    {'text': "Standard language for candidate parsing pipelines and machine learning development."}
                ]
            },
            {
                'runs': [
                    {'text': "Scikit-Learn (sklearn): ", 'bold': True},
                    {'text': "Powers `TfidfVectorizer` and `TruncatedSVD` (for Latent Semantic Analysis). Offers fast, robust, and reproducible vector models on standard CPU cores."}
                ]
            },
            {
                'runs': [
                    {'text': "NumPy & SciPy: ", 'bold': True},
                    {'text': "Optimized numerical operations for sparse matrix mathematics and cosine similarity calculations."}
                ]
            },
            {
                'runs': [
                    {'text': "Standard Libraries (csv, json, datetime): ", 'bold': True},
                    {'text': "Ensures lightweight execution without bulky external dependencies."}
                ]
            }
        ]
    },
    {
        'header': "Alternative Tech Rejections",
        'bullets': [
            {
                'runs': [
                    {'text': "No Deep Learning Models: ", 'bold': True},
                    {'text': "Transformers (e.g., Sentence-Transformers) require heavy download weights and slow down initialization on CPU servers."}
                ]
            },
            {
                'runs': [
                    {'text': "No Generative LLMs: ", 'bold': True},
                    {'text': "API calls for 100K profiles introduce severe cost, runtime latency, and non-deterministic ranking outcomes."}
                ]
            }
        ]
    }
]
format_slide_body(slide9, slide9_body, slide9_content)

# Slide 10: Submission Assets (index 9)
# Shape 2 (ID 119) is body
slide10 = prs.slides[9]
slide10_body = slide10.shapes[2]
slide10_content = [
    {
        'header': "Deliverable Package",
        'bullets': [
            {
                'runs': [
                    {'text': "GitHub Repository: ", 'bold': True},
                    {'text': "Includes modular source files, dependencies checklist, and a Colab playground notebook."}
                ]
            },
            {
                'runs': [
                    {'text': "submission.csv: ", 'bold': True},
                    {'text': "The final formatted output containing the top 100 candidates ranked by the hybrid engine."}
                ]
            },
            {
                'runs': [
                    {'text': "Approach PDF: ", 'bold': True},
                    {'text': "A converted copy of this presentation document explaining the system design."}
                ]
            }
        ]
    },
    {
        'header': "Validation & Schema Verification",
        'bullets': [
            {
                'runs': [
                    {'text': "Schema Verification: ", 'bold': True},
                    {'text': "The output structure matches `[candidate_id, rank, score, reasoning]` columns exactly."}
                ]
            },
            {
                'runs': [
                    {'text': "Automated Integrity Checks: ", 'bold': True},
                    {'text': "Validated via `validate_submission.py` to confirm that the file contains exactly 100 candidate rows, unique ranks from 1 to 100, scores are non-increasing by rank, and ties are broken in alphabetical order of IDs."}
                ]
            }
        ]
    }
]
format_slide_body(slide10, slide10_body, slide10_content)

# Save populated presentation
prs.save(output_path)
print(f"Presentation successfully saved to {output_path}")
